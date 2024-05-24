import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define a function to add a slide with title and content
def add_slide(title, content, bullet_points=None):
    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    # Set the title
    title_placeholder.text = title
    
    # Set the content
    tf = content_placeholder.text_frame
    tf.text = content
    
    # Add bullet points if provided
    if bullet_points:
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

# Add slides
add_slide("Meshduino Product Report", "")

add_slide("Introduction", 
          "Το MeshDuino είναι ένα καινοτόμο δίκτυο Arduino ενσωματωμένο στην Python",
          ["Διαισθητική γραφική διεπαφή χρήστη (GUI)",
           "Διευκόλυνση διαδραστικής μάθησης και πειραματισμού",
           "Δομημένες προκλήσεις και ανατροφοδότηση σε πραγματικό χρόνο"])

add_slide("Πως το Φανταζόμαστε", 
          "Δημιουργία εκπαιδευτικών παιχνιδιών STEM για παιδιά",
          ["Χρήση αισθητήρων και εξαρτημάτων",
           "Διδακτικά αντικείμενα: προγραμματισμός, περιβαλλοντική επιστήμη, φυσική κ.ά."])

add_slide("Στόχοι των Παιχνιδιών", 
          "",
          ["Παροχή διαδραστικής μαθησιακής εμπειρίας",
           "Εισαγωγή βασικών εννοιών προγραμματισμού",
           "Εκπαίδευση στην περιβαλλοντική παρακολούθηση και ηλεκτρονικά μέσω πειραμάτων"])

add_slide("Οδηγός GUI", 
          "Καθοδηγεί δασκάλους, γονείς, μαθητές ή ιδιώτες",
          ["Σεμινάριο για χρήση εφαρμογής και ρύθμιση Arduinos",
           "Απλά βήματα για εύκολη κατανόηση"])

add_slide("Προκλήσεις Εφαρμογής", 
          "Πολυάριθμες προκλήσεις σε διάφορους τομείς",
          ["Ηλεκτρολογία, Φυσική, Ακουστική, Τηλεπικοινωνίες, Μαθηματικά κ.ά.",
           "Διαδραστικές, διασκεδαστικές και εύκολες στη χρήση"])

add_slide("Βασικά Σημεία", 
          "",
          ["Προσέλκυση ενδιαφέροντος σε νεαρή ηλικία",
           "Ενθάρρυνση αναζήτησης προκλήσεων ή δημιουργίας έργων",
           "Πίνακας κατάταξης τοπικά, εθνικά και διεθνώς"])

add_slide("Δομή Εφαρμογής", 
          "Διαδραστικό GUI με διάφορες ενότητες",
          ["Home Section: Επισκόπηση έργου",
           "Challenge Section: Διαθέσιμες προκλήσεις",
           "Leaderboard Section: Κατατάξεις ομάδων",
           "About Us Section: Πληροφορίες ομάδας Meshduino"])

add_slide("Δομή Πρόκλησης", 
          "Κάθε πρόκληση περιλαμβάνει τρία πειράματα",
          ["Παράδειγμα: Engineering Challenge",
           "Διαδοχική εκτέλεση πειραμάτων με ανατροφοδότηση"])

add_slide("Πειράματα",
          "Πείραμα 1: Logic Gate LED Circuit",
          ["Εξαρτήματα: Arduino UNO, RFM22 Wireless Shield, LEDs, αντιστάσεις κ.ά.",
           "Περιγραφή: Συναρμολόγηση κυκλώματος λογικής πύλης"])

add_slide("Πειράματα (συνέχεια)",
          "Πείραμα 2: Code & Conquer - Virtual Maze Navigation",
          ["Εξαρτήματα: Arduino Uno, αισθητήρες υπερήχων, φωτοαντιστάσεις κ.ά.",
           "Περιγραφή: Προγραμματισμός για πλοήγηση σε εικονικό λαβύρινθο"])

add_slide("Πειράματα (συνέχεια)",
          "Πείραμα 3: Arduino Morse Code Encoder-Decoder",
          ["Εξαρτήματα: Arduino Uno, αντιστάσεις, κόκκινο LED, pushbutton",
           "Περιγραφή: Κωδικοποίηση και αποκωδικοποίηση μηνυμάτων σε κώδικα Morse"])

add_slide("Ολοκλήρωση Πρόκλησης", 
          "",
          ["Ενημέρωση GUI για εμφάνιση αποτελεσμάτων",
           "Καταχώρηση χρόνου και κατάταξης ομάδων"])

add_slide("Συμπεράσματα", 
          "Το Meshduino προσφέρει μια ολοκληρωμένη πλατφόρμα για μάθηση και πειραματισμό",
          ["Πρακτική εμπειρία με επεξεργασία δεδομένων, ενσωμάτωση αισθητήρων και συναρμολόγηση κυκλωμάτων",
           "Πολύτιμες δεξιότητες για μαθητές, ιδιώτες και γονείς"])

# Save the presentation
pptx_file = "/mnt/data/Meshduino_Product_Report.pptx"
prs.save(pptx_file)

pptx_file